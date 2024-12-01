from flask import Flask, request, send_file, jsonify
from pptx import Presentation
from pptx.util import Inches
import os
from pathlib import Path
from werkzeug.utils import secure_filename
import random
import string
from flask_cors import CORS

app = Flask(__name__)
CORS(app)

UPLOAD_FOLDER = 'uploads/'
OUTPUT_FOLDER = 'output/'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['OUTPUT_FOLDER'] = OUTPUT_FOLDER

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

if not os.path.exists(OUTPUT_FOLDER):
    os.makedirs(OUTPUT_FOLDER)

def generate_unique_filename(length=8):
    characters = string.ascii_lowercase + string.digits
    return ''.join(random.choice(characters) for i in range(length)) + '.pptx'

@app.route('/upload', methods=['POST'])
def upload_files():
    files = request.files.getlist('images')
    texts = request.form.getlist('texts')

    if not files or len(texts) == 0:
        return jsonify({"error": "No files or texts uploaded."}), 400

    prs = Presentation()

    img_start_idx = 0
    for i, text in enumerate(texts):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add all images for the current slide
        while img_start_idx < len(files):
            file = files[img_start_idx]
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            # Add picture to slide
            slide.shapes.add_picture(filepath, left=Inches(1), top=Inches(1), width=Inches(6))

            # Cleanup uploaded image
            os.remove(filepath)

            img_start_idx += 1
            # Break if we're moving to the next text
            if img_start_idx >= len(files) or (i + 1 < len(texts) and texts[i + 1]):
                break

        # Add text to slide
        text_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
        text_frame = text_box.text_frame
        text_frame.text = text

    pptx_filename = generate_unique_filename()
    pptx_filepath = os.path.join(app.config['OUTPUT_FOLDER'], pptx_filename)
    prs.save(pptx_filepath)

    return jsonify({"filename": pptx_filename})

@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    filepath = os.path.join(app.config['OUTPUT_FOLDER'], filename)
    if os.path.exists(filepath):
        return send_file(filepath, as_attachment=True)
    return jsonify({"error": "File not found"}), 404

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in {'png', 'jpg', 'jpeg', 'gif'}

if __name__ == '__main__':
    app.run(debug=True)
